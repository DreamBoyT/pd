import os  
import streamlit as st  
from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
from pptx.util import Inches  
from docx import Document  
from docx.shared import Inches as DocxInches  
from docx.oxml.ns import qn  
from docx.oxml import OxmlElement  
from io import BytesIO  
from langchain_openai import AzureChatOpenAI  
from langchain.prompts import PromptTemplate  
import re  
from PIL import Image  
import tempfile  
import pythoncom  
from win32com.client import Dispatch  
from pylatexenc.latex2text import LatexNodes2Text  
  
# Azure OpenAI API details  
azure_endpoint = 'https://ingitaraj-gpt4-turbo.openai.azure.com/openai/deployments/DanielIngitaraj-GPT4o/chat/completions?api-version=2024-02-15-preview'  
azure_deployment_name = 'DanielIngitaraj-GPT4o'  
azure_api_key = '3d84686d92be443e986d6215d7ce8f4b'  
azure_api_version = '2024-02-15-preview'  
  
# Initialize Azure OpenAI LLM  
llm = AzureChatOpenAI(  
    openai_api_key=azure_api_key,  
    api_version=azure_api_version,  
    azure_endpoint=azure_endpoint,  
    model="gpt-4o",  
    azure_deployment=azure_deployment_name,  
    temperature=0.3  
)  
  
# Profanity words list  
patent_profanity_words = [  
    "absolute", "absolutely", "all", "always", "authoritative", "authoritatively", "best", "biggest", "black hat",  
    "black list", "blackhat", "blacklist", "broadest", "certain", "certainly", "chinese wall", "compel", "compelled",  
    "compelling", "compulsorily", "compulsory", "conclusive", "conclusively", "constantly", "critical", "critically",  
    "crucial", "crucially", "decisive", "decisively", "definitely", "definitive", "definitively", "determinative",  
    "each", "earliest", "easiest", "embodiment", "embodiments", "entire", "entirely", "entirety", "essential",  
    "essentially", "essentials", "every", "everything", "everywhere", "exactly", "exclusive", "exclusively", "exemplary",  
    "exhaustive", "farthest", "finest", "foremost", "forever", "fundamental", "furthest", "greatest", "highest",  
    "imperative", "imperatively", "important", "importantly", "indispensable", "indispensably", "inescapable",  
    "inescapably", "inevitable", "inevitably", "inextricable", "inextricably", "inherent", "inherently", "instrumental",  
    "instrumentally", "integral", "integrally", "intrinsic", "intrinsically", "invaluable", "invaluably", "invariably",  
    "invention", "inventions", "irreplaceable", "irreplaceably", "key", "largest", "latest", "least", "littlest", "longest",  
    "lowest", "major", "man hours", "mandate", "mandated", "mandatorily", "mandatory", "master", "maximize", "maximum",  
    "minimize", "minimum", "most", "must", "nearest", "necessarily", "necessary", "necessitate", "necessitated",  
    "necessitates", "necessity", "need", "needed", "needs", "never", "newest", "nothing", "nowhere", "obvious", "obviously",  
    "oldest", "only", "optimal", "ought", "overarching", "paramount", "perfect", "perfected", "perfectly", "perpetual",  
    "perpetually", "pivotal", "pivotally", "poorest", "preferred", "purest", "required", "requirement", "requires",  
    "requisites", "shall", "shortest", "should", "simplest", "slaves", "slightest", "smallest", "tribal knowledge",  
    "ultimate", "ultimately", "unavoidable", "unavoidably", "unique", "uniquely", "unrivalled", "urgent", "urgently",  
    "valuable", "very", "vital", "vitally", "white hat", "white list", "whitehat", "whitelist", "widest", "worst"  
]  
  
# Lists for tone, style, and conditional & tentative language  
tone_list = [  
    "Precision and Specificity",  
    "Formality",  
    "Complexity",  
    "Objective and Impersonal",  
    "Structured and Systematic"  
]  
  
style_list = [  
    "Formal and Objective",  
    "Structured and Systematic",  
    "Technical Jargon and Terminology",  
    "Detailed and Specific",  
    "Impersonal Tone",  
    "Instructional and Descriptive",  
    "Use of Figures and Flowcharts",  
    "Legal and Protective Language",  
    "Repetitive and Redundant",  
    "Examples and Clauses"  
]  
  
conditional_language_list = [  
    "may include", "in some aspects", "aspects of the present disclosure", "wireless communication networks",  
    "by way of example", "may be", "may further include", "may be used", "may occur", "may use", "may monitor",  
    "may periodically wake up", "may demodulate", "may consume", "can be performed", "may enter and remain",  
    "may correspond to", "may also include", "may be identified in response to", "may be further a function of",  
    "may be multiplied by", "may schedule", "may select", "may also double", "may further comprise",  
    "may be configured to", "may correspond to a duration value", "may correspond to a product of", "may be closer",  
    "may be significant", "may not be able", "may result", "may reduce", "may be operating in", "may further be configured to",  
    "may further process", "may be executed by", "may be received", "may avoid", "may indicate", "may be selected",  
    "may be proactive", "may perform", "may be necessary", "may be amplified", "may involve", "may require", "may be stored",  
    "may be accessed", "may be transferred", "may be implemented", "may include instructions to", "may depend upon",  
    "may communicate", "may be generated", "may be configured"  
]  
  
# Function to sanitize text by removing non-XML-compatible characters  
def sanitize_text(text):  
    return re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)  
  
# Function to extract text and title from ppt slides  
def extract_text_and_title_from_ppt(ppt_file):  
    prs = Presentation(ppt_file)  
    slides_data = []  
    for slide_num, slide in enumerate(prs.slides, start=1):  
        slide_text = []  
        slide_title = None  
        for shape in slide.shapes:  
            if hasattr(shape, "text"):  
                if shape == slide.shapes.title:  
                    slide_title = shape.text  
                else:  
                    slide_text.append(shape.text)  
        slides_data.append((slide_num, slide_title, "\n".join(slide_text)))  
    return slides_data  
  
# Function to check if slide contains images, tables, or flowcharts  
def contains_relevant_elements(slide):  
    for shape in slide.shapes:  
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  
            return True  
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:  
            return True  
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and any(s.shape_type == MSO_SHAPE_TYPE.LINE for s in shape.shapes):  
            return True  
    return False  
  
# Function to remove bullet point text, title text, and header/footer contents from slides  
def remove_bullet_point_title_and_header_footer_text(prs):  
    for slide in prs.slides:  
        # Remove title text  
        if slide.shapes.title and hasattr(slide.shapes.title, "text_frame"):  
            slide.shapes.title.text_frame.clear()  
        # Remove bullet point text  
        for shape in slide.shapes:  
            if hasattr(shape, "text_frame") and shape.text_frame is not None:  
                for paragraph in shape.text_frame.paragraphs:  
                    if paragraph.text and (paragraph.level > 0 or paragraph.text.startswith(("\u2022", "\u2023", "\u25E6", "\u2043", "\u2219"))):  
                        paragraph.clear()  
        # Remove header and footer text  
        slide_headers_footers = slide.placeholders  
        for placeholder in slide_headers_footers:  
            if placeholder.is_placeholder and hasattr(placeholder, "text_frame"):  
                placeholder.text_frame.clear()  
    return prs  
  
# Function to remove logos and text from the slide master  
def remove_elements_from_master(prs):  
    for slide_master in prs.slide_masters:  
        for shape in slide_master.shapes:  
            if shape.has_text_frame:  
                shape.text_frame.clear()  
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  
                slide_master.shapes._spTree.remove(shape._element)  
    return prs  
  
# Function to identify potential logo dimensions and positions  
def identify_logo_shapes(prs, num_slides=5):  
    logo_shapes = {}  
    for slide in list(prs.slides)[:num_slides]:  
        for shape in slide.shapes:  
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  
                key = (shape.width, shape.height, shape.left, shape.top)  
                if key in logo_shapes:  
                    logo_shapes[key] += 1  
                else:  
                    logo_shapes[key] = 1  
  
    # Assume logos appear on most of the first few slides  
    logo_shapes = {k: v for k, v in logo_shapes.items() if v > num_slides / 2}  
    return logo_shapes  
  
# Function to remove logos from each slide based on identified dimensions and positions  
def remove_logos_from_slides(prs, logo_shapes):  
    for slide in prs.slides:  
        for shape in slide.shapes:  
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  
                key = (shape.width, shape.height, shape.left, shape.top)  
                if key in logo_shapes:  
                    slide.shapes._spTree.remove(shape._element)  
    return prs  
  
# Function to save slide as image using win32com  
def save_slide_as_image(prs, slide_num):  
    pythoncom.CoInitialize()  
    powerpoint = Dispatch("PowerPoint.Application")  
    powerpoint.Visible = 1  
  
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:  
        prs.save(tmp.name)  
        presentation = powerpoint.Presentations.Open(tmp.name)  
        slide = presentation.Slides(slide_num)  
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:  
            img_path = tmp_img.name  
        slide.Export(img_path, 'PNG')  
        presentation.Close()  
        powerpoint.Quit()  
    return img_path  
  
# Function to generate explanation using Azure OpenAI  
def generate_explanation(slide_text, figure_number):  
    prompt_template = PromptTemplate(  
        input_variables=["slide_text", "figure_number"],  
        template=f"""  
    Slide Content: {{slide_text}}  
      
    Aspects of the present disclosure may include insights extracted from the above slide content. The information should be delivered directly and engagingly in a single, coherent paragraph. Avoid phrases like 'The slide presents,' 'discusses,' 'outlines,' or 'content.' The explanation should be concise and semantically meaningful, summarizing all major points in one paragraph without line breaks or bullet points.  
      
    The text should adhere to the following style guidelines:  
    1. Remove all listed profanity words.  
    2. Use passive voice.  
    3. Use conditional and tentative language, such as "may include," "in some aspects," and "aspects of the present disclosure."  
    4. Replace "Million" with "1,000,000" and "Billion" with "1,000,000,000."  
    5. Maintain the following tone characteristics: {', '.join(tone_list)}.  
    6. Follow these style elements: {', '.join(style_list)}.  
    7. Use the following conditional and tentative language phrases: {', '.join(conditional_language_list)}.  
    8. Maintain the exact wording in the generated content. Do not substitute words with synonyms. For example, "instead" should remain "instead" and not be replaced with "conversely."  
    9. Replace the phrase "further development" with "our disclosure" in all generated content.  
    10. Make sure to use LaTeX formatting for all mathematical symbols, equations, subscripting, and superscripting to ensure they are displayed correctly in the output.  
    11. Reference the corresponding slide image as "Figure {{figure_number}}."  
    """  
    )  
    prompt = prompt_template.format(slide_text=sanitize_text(slide_text), figure_number=figure_number)  
    response = llm(prompt)  
    return response.content  
  
def add_math_symbols_to_doc(paragraph, text):  
    """  
    A function to ensure math symbols and equations are correctly added to a Word document.  
    """  
    parts = re.split(r'(\$[^\$]*\$)', text)  # Split by TeX math delimiters  
    for part in parts:  
        if part.startswith('$') and part.endswith('$'):  
            latex_text = part.strip('$')  
            # Convert LaTeX to OMML (Office Math Markup Language)  
            omml = LatexNodes2Text().latex_to_text(latex_text)  
            math_element = OxmlElement('m:oMathPara')  
            math_run = OxmlElement('m:oMath')  
            math_text = OxmlElement('m:t')  
            math_text.text = omml  
            math_run.append(math_text)  
            math_element.append(math_run)  
            paragraph._element.append(math_element)  
        else:  
            run = paragraph.add_run(part)  
    return paragraph  
  
# Streamlit app  
st.title("PPT Insights Extractor with Azure OpenAI")  
uploaded_file = st.file_uploader("Upload a PPT file", type=["pptx"])  
  
if uploaded_file is not None:  
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:  
        tmp.write(uploaded_file.read())  
        ppt_path = tmp.name  
  
    slides_data = extract_text_and_title_from_ppt(ppt_path)  
    prs = Presentation(ppt_path)  
    prs = remove_bullet_point_title_and_header_footer_text(prs)  
    prs = remove_elements_from_master(prs)  
  
    # Identify logo shapes from the first few slides  
    logo_shapes = identify_logo_shapes(prs, num_slides=5)  
    prs = remove_logos_from_slides(prs, logo_shapes)  
  
    explanations = []  
    for slide_num, slide_title, slide_text in slides_data:  
        explanation = generate_explanation(slide_text, slide_num)  
        explanations.append((slide_num, slide_title, explanation))  
  
    # Create a Word document  
    doc = Document()  
    style = doc.styles['Normal']  
    font = style.font  
    font.name = 'Times New Roman'  
    font.size = DocxInches(12 / 72)  # Convert points to inches  
    paragraph_format = style.paragraph_format  
    paragraph_format.line_spacing = 1.5  
    paragraph_format.alignment = 3  # Justify  
  
    for slide_num, slide_title, explanation in explanations:  
        sanitized_slide_title = sanitize_text(slide_title if slide_title else "Untitled Slide")  
        sanitized_explanation = sanitize_text(explanation)  
  
        # Add slide number and title inline with the content  
        paragraph = doc.add_paragraph()  
        run = paragraph.add_run(f"[{slide_num}, {sanitized_slide_title}] ")  
        # Add the explanation with math symbols handled  
        add_math_symbols_to_doc(paragraph, sanitized_explanation)  
  
    # Append slide images after the explanation, only for relevant slides  
    figure_counter = 1  
    for slide_num, _, _ in explanations:  
        slide = prs.slides[slide_num - 1]  
        if contains_relevant_elements(slide):  
            img_path = save_slide_as_image(prs, slide_num)  
            doc.add_paragraph(f"Figure {figure_counter}")  
            doc.add_picture(img_path, width=DocxInches(5))  
            os.remove(img_path)  # Clean up the image file  
            figure_counter += 1  
  
    os.remove(ppt_path)  # Clean up the temporary PPT file  
  
    # Save the Word document to a BytesIO object  
    buffer = BytesIO()  
    doc.save(buffer)  
    buffer.seek(0)  
  
    st.download_button(  
        label="Download Word Document",  
        data=buffer,  
        file_name="slides_insights.docx",  
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"  
    )
